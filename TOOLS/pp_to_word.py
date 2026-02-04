import os
import io
from pptx import Presentation
from docx import Document
from docx.shared import Inches
from PIL import Image

import re

def clean_xml_text(text: str) -> str:
    if not text:
        return ""
    # Remove invalid XML characters
    return re.sub(
        r"[\x00-\x08\x0B\x0C\x0E-\x1F]",
        "",
        text
    )

def pptx_to_word_auto_output(pptx_path: str):
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"File not found: {pptx_path}")

    # ---- Auto generate output path ----
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_dir = os.path.dirname(pptx_path)
    output_docx = os.path.join(output_dir, f"{base_name}_extracted.docx")

    prs = Presentation(pptx_path)
    doc = Document()

    doc.add_heading("PowerPoint Content Extraction", level=1)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        doc.add_heading(f"Slide {slide_idx}", level=2)

        # ---- Extract text ----
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text)

        if texts:
            doc.add_paragraph("Text content:")
            for t in texts:
                doc.add_paragraph(clean_xml_text(t), style="List Bullet")
        else:
            doc.add_paragraph("No text content.")

        # ---- Extract images ----
        img_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img_count += 1
                image_bytes = shape.image.blob
                image_stream = io.BytesIO(image_bytes)

                img = Image.open(image_stream)
                temp_img = f"_tmp_slide_{slide_idx}_{img_count}.png"
                img.save(temp_img)

                doc.add_paragraph(f"Image {img_count}:")
                doc.add_picture(temp_img, width=Inches(5))

                os.remove(temp_img)

        if img_count == 0:
            doc.add_paragraph("No images found.")

        doc.add_paragraph("—" * 30)

    doc.save(output_docx)
    print(f"✅ Word file created: {output_docx}")


if __name__ == "__main__":
    pptx_file = input("Enter PowerPoint file path: ").strip().strip("'").strip('"')
    pptx_to_word_auto_output(pptx_file)
